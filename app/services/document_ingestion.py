from __future__ import annotations

import re
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from uuid import UUID

from sqlmodel import Session, select

from app.core.config import settings
from app.core.errors import BlockedWorkflowError
from app.models.tables import UserSourceDocument, utc_now
from app.schemas.profile_document import SourceDocumentKind, SourceDocumentStatus
from app.services.url_normalizer import sha256_text

SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".docx", ".pptx"}
SUPPORTED_CONTENT_TYPES = {
    "text/plain",
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
MIN_USEFUL_TEXT_CHARS = 40


@dataclass(frozen=True)
class ExtractedDocumentText:
    text: str
    text_hash: str
    character_count: int
    parser_warnings: list[dict]


def ingest_source_document_bytes(
    session: Session,
    *,
    user_id: UUID,
    kind: SourceDocumentKind,
    filename: str,
    content_type: str | None,
    content: bytes,
) -> tuple[UserSourceDocument, ExtractedDocumentText, bool]:
    if not content:
        raise BlockedWorkflowError("uploaded source document is empty")
    if len(content) > settings.max_source_document_bytes:
        raise BlockedWorkflowError(
            "uploaded source document exceeds the configured size limit "
            f"({settings.max_source_document_bytes} bytes)"
        )

    safe_filename = _safe_filename(filename)
    extension = Path(safe_filename).suffix.lower()
    _validate_supported_file(extension, content_type)

    file_hash = _sha256_bytes(content)
    existing = session.exec(
        select(UserSourceDocument).where(
            UserSourceDocument.user_id == user_id,
            UserSourceDocument.file_hash == file_hash,
        )
    ).first()

    reused_existing = existing is not None
    document = existing or UserSourceDocument(
        user_id=user_id,
        kind=kind.value,
        filename=safe_filename,
        content_type=content_type,
        file_hash=file_hash,
        storage_path="",
        status=SourceDocumentStatus.UPLOADED.value,
    )
    if not existing:
        session.add(document)
        session.flush()

    storage_path = _document_storage_path(
        user_id=user_id,
        document_id=document.id,
        filename=safe_filename,
    )
    storage_path.parent.mkdir(parents=True, exist_ok=True)
    if not storage_path.exists():
        storage_path.write_bytes(content)

    extracted = extract_text_from_document(
        filename=safe_filename,
        content_type=content_type,
        content=content,
    )
    document.kind = kind.value
    document.filename = safe_filename
    document.content_type = content_type
    document.storage_path = str(storage_path)
    document.extracted_text_hash = extracted.text_hash
    document.character_count = extracted.character_count
    document.parser_warnings = extracted.parser_warnings
    document.status = SourceDocumentStatus.TEXT_EXTRACTED.value
    document.error_message = None
    document.updated_at = utc_now()
    session.flush()
    return document, extracted, reused_existing


def extract_text_from_document(
    *,
    filename: str,
    content_type: str | None,
    content: bytes,
) -> ExtractedDocumentText:
    extension = Path(filename).suffix.lower()
    warnings: list[dict] = []
    try:
        if extension == ".txt":
            text = _extract_txt(content)
        elif extension == ".pdf":
            text, parser_warnings = _extract_pdf(content)
            warnings.extend(parser_warnings)
        elif extension == ".docx":
            text, parser_warnings = _extract_docx(content)
            warnings.extend(parser_warnings)
        elif extension == ".pptx":
            text, parser_warnings = _extract_pptx(content)
            warnings.extend(parser_warnings)
        else:
            raise BlockedWorkflowError(f"unsupported source document extension: {extension}")
    except BlockedWorkflowError:
        raise
    except Exception as exc:
        raise BlockedWorkflowError(f"source document text extraction failed: {exc}") from exc

    normalized = _normalize_document_text(text)
    if len(normalized) < MIN_USEFUL_TEXT_CHARS:
        raise BlockedWorkflowError(
            "source document does not contain enough readable text for profile extraction"
        )
    if len(normalized) > settings.max_source_document_text_chars:
        warnings.append(
            {
                "code": "text_truncated",
                "message": (
                    "Extracted source document text exceeded the configured analysis limit; "
                    "only the first configured characters were analyzed."
                ),
                "severity": "medium",
            }
        )
        normalized = normalized[: settings.max_source_document_text_chars]

    return ExtractedDocumentText(
        text=normalized,
        text_hash=sha256_text(normalized),
        character_count=len(normalized),
        parser_warnings=warnings,
    )


def mark_document_failed(
    session: Session,
    *,
    document: UserSourceDocument,
    error_message: str,
) -> None:
    document.status = SourceDocumentStatus.FAILED.value
    document.error_message = error_message
    document.updated_at = utc_now()
    session.add(document)


def _extract_txt(content: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise BlockedWorkflowError("text source document could not be decoded")


def _extract_pdf(content: bytes) -> tuple[str, list[dict]]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise BlockedWorkflowError("PDF extraction requires the pypdf package") from exc

    reader = PdfReader(BytesIO(content))
    if getattr(reader, "is_encrypted", False):
        raise BlockedWorkflowError("password-protected PDF files are not supported")

    warnings: list[dict] = []
    pages: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception:
            warnings.append(
                {
                    "code": "pdf_page_unreadable",
                    "message": f"Could not extract text from PDF page {index}.",
                    "severity": "medium",
                }
            )
            page_text = ""
        if page_text.strip():
            pages.append(f"[page {index}]\n{page_text}")

    if not pages:
        raise BlockedWorkflowError(
            "PDF appears to be scanned or image-only; upload a text-based PDF/DOCX or paste text"
        )
    return "\n\n".join(pages), warnings


def _extract_docx(content: bytes) -> tuple[str, list[dict]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise BlockedWorkflowError("DOCX extraction requires the python-docx package") from exc

    document = Document(BytesIO(content))
    chunks: list[str] = []
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            chunks.append(paragraph.text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                chunks.append(" | ".join(cells))
    return "\n".join(chunks), []


def _extract_pptx(content: bytes) -> tuple[str, list[dict]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise BlockedWorkflowError("PPTX extraction requires the python-pptx package") from exc

    presentation = Presentation(BytesIO(content))
    chunks: list[str] = []
    warnings: list[dict] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_chunks: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_chunks.append(text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                slide_chunks.append(f"Speaker notes: {notes}")
        if slide_chunks:
            chunks.append(f"[slide {slide_index}]\n" + "\n".join(slide_chunks))
        else:
            warnings.append(
                {
                    "code": "pptx_slide_without_text",
                    "message": f"Slide {slide_index} did not contain extractable text.",
                    "severity": "info",
                }
            )
    return "\n\n".join(chunks), warnings


def _validate_supported_file(extension: str, content_type: str | None) -> None:
    if extension not in SUPPORTED_EXTENSIONS:
        raise BlockedWorkflowError(
            "unsupported source document type; upload PDF, DOCX, PPTX, or TXT"
        )
    # Browsers occasionally send octet-stream, so extension remains the source of truth.
    if (
        content_type
        and content_type not in SUPPORTED_CONTENT_TYPES
        and content_type != "application/octet-stream"
    ):
        raise BlockedWorkflowError(
            "unsupported source document content type; upload PDF, DOCX, PPTX, or TXT"
        )


def _normalize_document_text(text: str) -> str:
    lines: list[str] = []
    for raw_line in text.replace("\x00", "").splitlines():
        line = " ".join(raw_line.split())
        if line:
            lines.append(line)
    return "\n".join(lines)


def _safe_filename(filename: str) -> str:
    name = Path(filename or "source-document").name.strip() or "source-document"
    return re.sub(r"[^A-Za-z0-9._ -]+", "_", name)[:180]


def _document_storage_path(*, user_id: UUID, document_id: UUID, filename: str) -> Path:
    return settings.source_documents_dir / str(user_id) / str(document_id) / filename


def _sha256_bytes(content: bytes) -> str:
    import hashlib

    return hashlib.sha256(content).hexdigest()
